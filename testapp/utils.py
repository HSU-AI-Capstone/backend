import logging
import os
import shutil
import subprocess
import sys
import tempfile
import uuid
from pathlib import Path

from .create_ppt import build_lecture_video
from .pdf2text import extract_text_from_pdf_content
from .prompts import ppt_gen_prompt
from .use_gpt import (
    generate_lesson_script,
    MODEL_NAME,
    get_openai_client,
    clean_text_with_llm,
    API_KEY,
)
from .voice import tts_pages_to_mp3, DEFAULT_VOICE_KEY

logger = logging.getLogger(__name__)


def mock_generate_lecture_video(pdf_path: str) -> str:
    """
    실제 영상 생성 대신 testapp/static/dummy.mp4를 복사해서 사용하는 테스트용 함수.
    """
    base_dir = os.path.dirname(os.path.dirname(__file__))  # backend 기준에서 한 단계 위
    dummy_video_path = os.path.join(
        base_dir, "testapp", "static", "final_lecture_video.mp4"
    )

    if not os.path.exists(dummy_video_path):
        raise FileNotFoundError(
            "해당 mp4 파일이 존재하지 않습니다: " + dummy_video_path
        )

    # 임시 디렉토리에 복사
    temp_dir = tempfile.mkdtemp()
    target_path = os.path.join(temp_dir, "lecture_video.mp4")
    shutil.copyfile(dummy_video_path, target_path)

    return target_path


def generate_lecture_video(
    subject: str,
    description: str,
    professor: str,
    pdf_path: str,
) -> str:
    """
    사용자의 입력(subject, description, professor, pdf_path)을 받아
    AI로 PPTX, 대본, 오디오를 순차 생성하고 마지막에 MP4 비디오 경로를 반환한다.
    """
    # 1) 작업 디렉터리 생성 → 모든 중간 산출물(tmpdir) 자동 삭제
    temp_dir = tempfile.mkdtemp(prefix="lecture_gen_")
    try:
        workdir = Path(temp_dir)

        # ───────────────────────────────────────────────
        # 2) PDF → 텍스트 → (옵션) 정제
        # ───────────────────────────────────────────────
        pdf_bytes = Path(pdf_path).read_bytes()
        raw_text = extract_text_from_pdf_content(pdf_bytes)
        if raw_text is None:
            raise RuntimeError("PDF 텍스트를 추출하지 못했습니다.")
        # 필요시 LLM으로 노이즈 제거
        cleaned = clean_text_with_llm(raw_text, API_KEY, MODEL_NAME) or raw_text
        text_file = workdir / "lecture_text.txt"
        text_file.write_text(cleaned, encoding="utf-8")

        # ───────────────────────────────────────────────
        # 3) LLM으로 PPTX 생성 코드 받아 실행
        # ───────────────────────────────────────────────
        # 3.1 LLM 클라이언트 초기화
        client = get_openai_client(API_KEY)
        if client is None:
            raise RuntimeError("OpenAI 클라이언트 초기화 실패")

        # 3.2 Python-pptx 생성용 프롬프트 구성
        prompt = ppt_gen_prompt + "\n\n" + cleaned

        # 3.3 코드 생성 요청
        response = client.chat.completions.create(
            model=MODEL_NAME,
            messages=[
                {
                    "role": "system",
                    "content": "당신은 python-pptx 코드를 생성하는 AI입니다. 유효한 Python 코드만 리턴하세요.",
                },
                {"role": "user", "content": prompt},
            ],
            temperature=0.3,
        )
        raw_code = response.choices[0].message.content

        # 3.4 코드 블록 마커 제거
        def strip_fences(raw: str) -> str:
            """
            ```python ... ``` 또는 ``` ... ``` 로 감싸인 코드를
            언어 식별자까지 포함해 깔끔히 추출합니다.
            """
            s = raw.strip()

            # 1) ``` 로 분할하고 중간 부분 취득
            if "```" in s:
                parts = s.split("```")
                # parts = ["", "python\n<code...>\n", ...]
                content = parts[1]
            else:
                content = s

            # 2) 만약 첫 줄이 'python' 이라면 제거
            lines = content.splitlines()
            if lines and lines[0].strip().lower() == "python":
                lines = lines[1:]

            # 3) 끝에 ``` 가 남아 있다면 제거
            if lines and lines[-1].strip() == "```":
                lines = lines[:-1]

            return "\n".join(lines).strip()

        # 3.4 코드 블록 마커·언어 식별자 제거
        code = strip_fences(raw_code)

        # 3.4.1 데이터 구조 검증 코드 추가
        validation_code = """
def validate_slide_data(slides_data):
    \"\"\"슬라이드 데이터 구조를 검증합니다.\"\"\"
    if not isinstance(slides_data, list):
        raise ValueError("slides_data는 리스트여야 합니다.")
    
    for slide in slides_data:
        if not isinstance(slide, dict):
            raise ValueError("각 슬라이드는 딕셔너리여야 합니다.")
        
        if "points" not in slide:
            raise ValueError("각 슬라이드는 'points' 키를 포함해야 합니다.")
        
        if not isinstance(slide["points"], list):
            raise ValueError("points는 리스트여야 합니다.")
        
        for point in slide["points"]:
            if not isinstance(point, dict):
                raise ValueError("각 point는 딕셔너리여야 합니다.")
            if "text" not in point:
                raise ValueError("각 point는 'text' 키를 포함해야 합니다.")
"""

        # 검증 코드를 structured_slides 정의 직후에 삽입
        if "structured_slides =" in code:
            # structured_slides 정의 후에 검증 코드 삽입
            code = code.replace(
                "structured_slides =",
                "structured_slides =\n\n"
                + validation_code
                + "\n# 데이터 검증\nvalidate_slide_data(structured_slides)\n\n",
            )

        # 3.5 파일로 저장 후 실행
        code_file = workdir / "gen_ppt.py"
        code_file.write_text(code, encoding="utf-8")
        print(f"PPTX 생성 코드 저장 완료: {code_file}")

        try:
            # 코드 실행
            result = subprocess.run(
                [sys.executable, str(code_file)],
                check=True,
                cwd=str(workdir),
                capture_output=True,
                text=True,
            )
            print(f"PPTX 생성 코드 실행 결과:")
            print(f"stdout: {result.stdout}")
            print(f"stderr: {result.stderr}")
        except subprocess.CalledProcessError as e:
            print(f"PPTX 생성 코드 실행 실패:")
            print(f"stdout: {e.stdout}")
            print(f"stderr: {e.stderr}")
            raise RuntimeError(f"PPTX 생성 코드 실행 실패: {str(e)}")

        # 3.6 생성된 .pptx 파일 찾기 및 검증
        pptx_list = list(workdir.glob("*.pptx"))
        print(f"찾은 PPTX 파일들: {pptx_list}")

        if not pptx_list:
            raise RuntimeError("PPTX 파일이 생성되지 않았습니다.")

        pptx_path = str(pptx_list[0])
        print(f"사용할 PPTX 파일: {pptx_path}")

        # PPTX 파일 크기 확인
        pptx_size = os.path.getsize(pptx_path)
        if pptx_size == 0:
            raise RuntimeError("생성된 PPTX 파일이 비어있습니다.")
        print(f"PPTX 파일 크기: {pptx_size} bytes")

        # PPTX 구조 추출
        from pptx import Presentation

        prs = Presentation(pptx_path)
        ppt_structure = []
        for slide in prs.slides:
            slide_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_content.append(shape.text)
            ppt_structure.append("\n".join(slide_content))

        # ───────────────────────────────────────────────
        # 4) 대본 생성 → 페이지별 MP3 변환
        # ───────────────────────────────────────────────
        script_file = workdir / "lesson_script.txt"
        lesson = generate_lesson_script(
            input_text_file=str(text_file),
            output_script_file=str(script_file),
            description=description,
            model=MODEL_NAME,
            ppt_structure=ppt_structure,  # PPT 구조 전달
        )
        if lesson is None:
            raise RuntimeError("수업 대본 생성 실패")

        audio_dir = workdir / "audio"
        # 기본 목소리: DEFAULT_VOICE_KEY ("DAWOON" 등)
        tts_pages_to_mp3(
            txt_path=str(script_file),
            out_dir=str(audio_dir),
            voice_key=DEFAULT_VOICE_KEY,
            base_name="page",
        )

        # ───────────────────────────────────────────────
        # 5) 슬라이드 + 오디오 합성 → MP4 생성
        # ───────────────────────────────────────────────
        video_filename = f"{uuid.uuid4().hex}.mp4"
        video_path = workdir / video_filename
        build_lecture_video(
            pptx_file=pptx_path,
            audio_dir=str(audio_dir),
            output_path=str(video_path),
            fps=24,
        )

        # 비디오 파일이 생성되었는지 확인
        if not video_path.exists():
            raise RuntimeError(f"비디오 파일이 생성되지 않았습니다: {video_path}")

        # 비디오 파일 크기 확인
        video_size = os.path.getsize(video_path)
        if video_size == 0:
            raise RuntimeError(f"생성된 비디오 파일이 비어있습니다: {video_path}")
        print(f"비디오 파일 크기: {video_size} bytes")

        return str(video_path)

    except Exception as e:
        # 오류 발생 시 임시 디렉토리 정리
        shutil.rmtree(temp_dir, ignore_errors=True)
        raise e
